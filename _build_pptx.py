"""One-shot generator for Thesis_presentation_v2.pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


FONT = "Times New Roman"
DARK_BLUE = RGBColor(0x00, 0x20, 0x60)
ACCENT_BLUE = RGBColor(0x32, 0x5A, 0xDA)
TEAL = RGBColor(0x0C, 0xA0, 0xA4)
GREEN = RGBColor(0x00, 0xB0, 0x50)
LIGHT_BLUE = RGBColor(0x00, 0xB0, 0xF0)
GREY = RGBColor(0x59, 0x59, 0x59)
BG = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xFA)


def add_title_bar(slide, text, color=DARK_BLUE, size=28):
    """Title strip at the top of a content slide."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    # decorative underline
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.05), Inches(2.0), Inches(1.05))
    line.line.color.rgb = ACCENT_BLUE
    line.line.width = Pt(3)
    return box


def add_textbox(slide, left, top, width, height, items, *, title=None,
                title_color=ACCENT_BLUE, body_size=14, title_size=16,
                border=True, bullet=True, body_color=None):
    """Add a text box with optional title + bullet list."""
    box = slide.shapes.add_textbox(left, top, width, height)
    if border:
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    first = True
    if title is not None:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = FONT
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False

    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = (("• " + item) if bullet else item)
        run.font.name = FONT
        run.font.size = Pt(body_size)
        if body_color is not None:
            run.font.color.rgb = body_color
    return box


def add_table(slide, left, top, width, height, headers, rows, *,
              header_color=ACCENT_BLUE, font_size=12, header_size=12,
              highlight_cell=None, highlight_color=GREEN):
    """Add a styled table. `highlight_cell` = (row_idx, col_idx) (0-based, excluding header)."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = FONT
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = BG

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT
            run.font.size = Pt(font_size)
            if highlight_cell == (i - 1, j):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE2, 0xEF, 0xDA)
                run.font.bold = True
                run.font.color.rgb = highlight_color
    return tbl_shape


def add_numbered_circle(slide, left, top, n, color):
    """Small numbered circle (used in slide 2)."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(n)
    run.font.name = FONT
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BG


# -----------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # -----------------------------------------------------------------
    # Slide 1 — Title
    s = prs.slides.add_slide(blank)
    # header
    box = s.shapes.add_textbox(Inches(2.0), Inches(0.3), Inches(9.3), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("Министерство науки и высшего образования Российской Федерации\n"
                "Челябинский государственный университет")
    run.font.name = FONT
    run.font.size = Pt(14)

    # title
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("Разработка программного комплекса для\n"
                "детектирования дрифта данных и адаптивного\n"
                "переобучения моделей машинного обучения")
    run.font.name = FONT
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    # student info
    box = s.shapes.add_textbox(Inches(8.0), Inches(6.0), Inches(5.0), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Студент: Н. М. Джордана"
    run.font.name = FONT
    run.font.size = Pt(18)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "группа МПрИ-201"
    run.font.name = FONT
    run.font.size = Pt(18)

    # -----------------------------------------------------------------
    # Slide 2 — Проблема
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Проблема: data drift и риски в продакшене")

    cards = [
        ("Data drift", "Сдвиг распределений входных данных относительно "
         "распределений при обучении.", ACCENT_BLUE),
        ("Скрытый риск", "Снижение качества модели «тихо», без явных "
         "ошибок в коде сервиса.", LIGHT_BLUE),
        ("Недостаточность детекторов", "Нужен аудит, переобучение, "
         "версионирование и деплой, а не просто алерт.", TEAL),
        ("Фрагментация инструментов", "Evidently / NannyML / MLflow / "
         "Airflow покрывают разные части — нужна интеграционная связка.",
         GREEN),
    ]
    top_y = 1.4
    for i, (title, body, color) in enumerate(cards):
        y = top_y + i * 1.35
        add_numbered_circle(s, Inches(0.9), Inches(y), i + 1, color)
        add_textbox(
            s, Inches(1.6), Inches(y - 0.05), Inches(11.2), Inches(1.1),
            [body], title=title, title_color=color, body_size=15, title_size=16,
            border=True, bullet=False,
        )

    # bottom highlight
    box = s.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("Цель работы — единая платформа вместо набора "
                "разрозненных MLOps-инструментов.")
    run.font.name = FONT
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    # -----------------------------------------------------------------
    # Slide 3 — Цель и задачи
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Цель и задачи работы")

    add_textbox(
        s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.5),
        ["Разработка программного комплекса для детектирования дрифта данных "
         "и адаптивного переобучения моделей машинного обучения."],
        title="Цель работы", title_color=ACCENT_BLUE, body_size=18, title_size=18,
        border=True, bullet=False,
    )

    tasks = [
        "Выполнить анализ предметной области.",
        "Спроектировать архитектуру интегрированной MLOps-платформы.",
        "Реализовать механизмы детектирования дрифта, оркестрации, автоматического "
        "переобучения и управления жизненным циклом моделей.",
        "Разработать REST API и веб-панель мониторинга.",
        "Провести экспериментальную валидацию на реальных и синтетических наборах данных.",
        "Сформулировать выводы и направления дальнейшего развития.",
    ]
    add_textbox(
        s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(4.0),
        [f"{i+1}. {t}" for i, t in enumerate(tasks)],
        title="Задачи", title_color=ACCENT_BLUE, body_size=16, title_size=18,
        border=True, bullet=False,
    )

    # -----------------------------------------------------------------
    # Slide 4 — Сравнение аналогов
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Сравнение с существующими инструментами")

    headers = ["Инструмент", "Дрифт", "Реестр\nмоделей", "Оркестрация",
               "Автоперео-\nбучение", "Open-source"]
    rows = [
        ["Evidently AI",       "●", "○", "○", "○", "●"],
        ["NannyML",            "●", "○", "○", "○", "●"],
        ["MLflow",             "◐", "●", "○", "○", "●"],
        ["Weights & Biases",   "◐", "●", "○", "○", "◐"],
        ["Apache Airflow",     "○", "○", "●", "◐", "●"],
        ["Fiddler AI",         "●", "◐", "○", "○", "○"],
        ["Наша платформа",     "●", "●", "●", "●", "●"],
    ]
    add_table(
        s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(4.0),
        headers, rows, font_size=14, header_size=14,
        highlight_cell=(6, 0),  # "Наша платформа" cell
    )

    add_textbox(
        s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.4),
        ["Ни одно из существующих решений не покрывает полный цикл "
         "drift → policy → retrain → lifecycle → inference в одной системе."],
        title="Вывод", title_color=DARK_BLUE, body_size=16, title_size=16,
        border=True, bullet=False,
    )

    # -----------------------------------------------------------------
    # Slide 5 — Статистические методы
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Статистические методы детектирования дрифта")

    methods = [
        ("PSI (Population Stability Index)", ACCENT_BLUE, [
            "Формула: PSI = Σ (Qᵢ − Pᵢ) · ln(Qᵢ / Pᵢ)",
            "Бины: 10 квантильных интервалов",
            "Бэнды: < 0.10 stable; 0.10–0.25 moderate; ≥ 0.25 high",
            "Интерпретируемость для стейкхолдеров",
        ]),
        ("KS (Колмогоров–Смирнов)", LIGHT_BLUE, [
            "Статистика: D = supₓ |Fₙ(x) − Gₘ(x)|",
            "α = 0.05",
            "Численные признаки",
            "Нон-параметрический, дистрибутив-фри",
        ]),
        ("χ² (Хи-квадрат)", TEAL, [
            "Статистика: χ² = Σ (Oᵢ − Eᵢ)² / Eᵢ",
            "α = 0.05",
            "Категориальные признаки",
            "Объединение редких категорий в `__other__`",
        ]),
    ]
    col_w = 4.0
    for i, (title, color, body) in enumerate(methods):
        x = 0.5 + i * (col_w + 0.15)
        add_textbox(
            s, Inches(x), Inches(1.4), Inches(col_w), Inches(4.0),
            body, title=title, title_color=color, body_size=14, title_size=15,
            border=True, bullet=True,
        )

    add_textbox(
        s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2),
        ["Три метода покрывают численные + категориальные признаки, с низкой "
         "вычислительной сложностью и интерпретируемым выходом."],
        title="Сводка", title_color=DARK_BLUE, body_size=14, title_size=15,
        border=True, bullet=False,
    )

    # -----------------------------------------------------------------
    # Slide 6 — Архитектура
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Архитектура системы: 5 модулей + 2 интерфейса")

    # Top: interfaces
    add_textbox(
        s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.7),
        ["REST API (FastAPI)         •         SPA Dashboard (React / MUI)"],
        title=None, body_size=18, border=True, bullet=False,
    )

    # 5 modules in a row
    modules = [
        ("drift_detection", "PSI / KS / χ²,\nBaselineProfile,\nDriftReport", ACCENT_BLUE),
        ("orchestration", "Engine + Policy +\nActions + RunStore", ACCENT_BLUE),
        ("retraining", "Pipeline + Registry\n+ Promotion rule", ACCENT_BLUE),
        ("lifecycle", "Stages + experiments\n+ model versions", ACCENT_BLUE),
        ("data_management", "Dataset versions +\nbaselines + provenance", ACCENT_BLUE),
    ]
    mod_w = 2.35
    for i, (name, body, color) in enumerate(modules):
        x = 0.5 + i * (mod_w + 0.1)
        add_textbox(
            s, Inches(x), Inches(2.4), Inches(mod_w), Inches(2.5),
            [body], title=name, title_color=color, body_size=13, title_size=14,
            border=True, bullet=False,
        )

    # Storage layer
    add_textbox(
        s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.4),
        ["SQLite (3 БД): orchestration.db • lifecycle.db • data_management.db",
         "Артефакты моделей: artifacts/models/model_vN.joblib + registry.json + champion.json"],
        title="Хранение", title_color=DARK_BLUE, body_size=14, title_size=15,
        border=True, bullet=True,
    )

    # Caption
    box = s.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("Один запуск оркестрации = drift report → policy → actions → "
                "(optional) retrain → lifecycle update → provenance.")
    run.font.name = FONT
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = GREY

    # -----------------------------------------------------------------
    # Slide 7 — Политика и пороги
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Политика и пороги (DriftThresholdPolicy)")

    headers = ["Счётчик", "Порог", "Логика"]
    rows = [
        ["n_features_high_psi", "> 0",
         "Любая «high»-PSI фича запускает триггер"],
        ["n_numeric_ks_significant", "> 2",
         "Допускаем ~2 KS-FP при α = 0.05"],
        ["n_categorical_chi2_significant", "> 3",
         "Допускаем ~3 χ²-FP при α = 0.05"],
    ]
    add_table(
        s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.0),
        headers, rows, font_size=14, header_size=14,
    )

    actions = [
        "LogAction — структурированное логирование причин срабатывания.",
        "WebhookAlertAction — опциональный POST на внешний webhook.",
        "RetrainPipelineAction — запуск пайплайна автоматического переобучения.",
    ]
    add_textbox(
        s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(2.0),
        actions, title="Action chain (выполняется при срабатывании)",
        title_color=ACCENT_BLUE, body_size=14, title_size=15,
        border=True, bullet=True,
    )

    add_textbox(
        s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.4),
        ["Каждый запуск (триггерный или нет) сохраняется в SQLite "
         "(artifacts/orchestration.db) с timestamp, scenario, причинами "
         "и JSON-сводкой отчёта о дрейфе."],
        title="Аудит", title_color=DARK_BLUE, body_size=14, title_size=15,
        border=True, bullet=False,
    )

    # -----------------------------------------------------------------
    # Slide 8 — Жизненный цикл моделей
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Жизненный цикл моделей")

    # Stage cards in a horizontal flow
    stages = [
        ("development", ACCENT_BLUE),
        ("staging", LIGHT_BLUE),
        ("production", GREEN),
        ("archived", GREY),
    ]
    card_w = 2.4
    gap = 0.3
    start_x = 0.8
    for i, (name, color) in enumerate(stages):
        x = start_x + i * (card_w + gap)
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(1.6),